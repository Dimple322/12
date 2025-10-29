from pathlib import Path
from typing import Dict
import re
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from agent.utils.column_normalizer import normalize_dataframe_columns

def extract_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _pdf_text(file_path)
    if suffix in {".xlsx", ".xls"}:
        return _excel_text(file_path)
    if suffix == ".pptx":
        return _pptx_text(file_path)
    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"Неизвестное расширение: {suffix}")

def _pdf_text(path: Path) -> str:
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def _excel_text(path: Path) -> str:
    df = pd.read_excel(path, sheet_name=0, dtype=str).fillna("")

    # Нормализация имён столбцов через общую функцию
    df_normalized = normalize_dataframe_columns(df)

    # Отладка (опционально)
    # print(f"[DEBUG] Нормализованные имена столбцов в _excel_text: {list(df_normalized.columns)}")

    return df_normalized.to_csv(index=False)

def _pptx_text(path: Path) -> str:
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
    return "\n".join(lines)